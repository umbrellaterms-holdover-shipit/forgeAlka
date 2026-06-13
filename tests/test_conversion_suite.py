from pathlib import Path
import csv
import json

from apexdev.cli import main
from apexdev.documents.conversion import convert_file, available_conversions
from apexdev.documents.docx_tools import docx_to_markdown
from apexdev.documents.pdf_tools import pdf_to_text
from apexdev.system_tools import dependency_report, apt_install_command


def test_available_conversions_include_core_pairs():
    pairs = {(row["from"], row["to"]) for row in available_conversions()}
    assert ("md", "docx") in pairs
    assert ("docx", "md") in pairs
    assert ("csv", "xlsx") in pairs
    assert ("xlsx", "csv") in pairs
    assert ("pptx", "md") in pairs


def test_markdown_docx_roundtrip_and_cli(tmp_path):
    md = tmp_path / "in.md"
    md.write_text("# Title\n\n- one\n- two\n\n| A | B |\n| --- | --- |\n| 1 | 2 |\n", encoding="utf-8")
    docx = tmp_path / "out.docx"
    result = convert_file(md, docx)
    assert result.method == "python-docx"
    assert docx.exists()
    extracted = docx_to_markdown(docx)
    assert "# Title" in extracted
    assert "| A | B |" in extracted

    back = tmp_path / "back.md"
    main(["convert", str(docx), str(back), "--json"])
    assert "# Title" in back.read_text(encoding="utf-8")


def test_markdown_pdf_and_pdf_text(tmp_path):
    md = tmp_path / "paper.md"
    md.write_text("# PDF Title\n\nSome body text.", encoding="utf-8")
    pdf = tmp_path / "paper.pdf"
    result = convert_file(md, pdf, prefer_pandoc=False)
    assert result.method == "reportlab-fallback"
    assert pdf.exists() and pdf.stat().st_size > 100
    txt = tmp_path / "paper.txt"
    result = convert_file(pdf, txt)
    assert result.method == "pypdf"
    assert "PDF Title" in txt.read_text(encoding="utf-8")


def test_csv_xlsx_roundtrip(tmp_path):
    csv_path = tmp_path / "rows.csv"
    csv_path.write_text("name,value\na,1\nb,2\n", encoding="utf-8")
    xlsx = tmp_path / "rows.xlsx"
    result = convert_file(csv_path, xlsx)
    assert result.method == "openpyxl"
    back = tmp_path / "back.csv"
    convert_file(xlsx, back)
    rows = list(csv.reader(back.open(encoding="utf-8")))
    assert rows[0] == ["name", "value"]
    assert rows[1] == ["a", "1"]


def test_pptx_to_markdown(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Deck Title"
    slide.placeholders[1].text = "Bullet body"
    pptx = tmp_path / "deck.pptx"
    prs.save(str(pptx))

    md = tmp_path / "deck.md"
    result = convert_file(pptx, md)
    assert result.method == "python-pptx"
    text = md.read_text(encoding="utf-8")
    assert "Deck Title" in text
    assert "Bullet body" in text


def test_deps_doctor_and_dry_run_cli(tmp_path, capsys):
    statuses = dependency_report(["ffmpeg", "pandoc"])
    assert {s.name for s in statuses} == {"ffmpeg", "pandoc"}
    main(["deps", "doctor"])
    captured = capsys.readouterr().out
    assert "Dependency Report" in captured
    main(["deps", "install", "ffmpeg", "pandoc", "--dry-run", "--json"])
    captured = capsys.readouterr().out
    assert captured.strip() == "[]" or "apt-get" in captured
