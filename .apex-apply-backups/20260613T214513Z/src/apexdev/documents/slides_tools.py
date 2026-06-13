"""PowerPoint slide generation/extraction utilities using python-pptx."""

from __future__ import annotations

from pathlib import Path

from ..core.models import Conversation

try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Pt  # type: ignore
except ImportError as exc:  # pragma: no cover - optional dependency
    Presentation = None  # type: ignore
    _pptx_import_error = exc
else:
    _pptx_import_error = None


def _require_pptx():
    if Presentation is None:
        raise ImportError("python-pptx is required for slide conversion but is not installed") from _pptx_import_error


def conversation_to_pptx(conversation: Conversation, output_path: str | Path) -> None:
    _require_pptx()
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = conversation.title or 'Conversation'
    if conversation.messages:
        subtitle = slide.placeholders[1]
        subtitle.text = f"{len(conversation.messages)} messages"
    bullet_slide_layout = prs.slide_layouts[1]
    for msg in conversation.messages:
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_placeholder = slide.shapes.title
        body = slide.placeholders[1]
        role = msg.role.capitalize() if msg.role else 'Message'
        title_placeholder.text = role
        tf = body.text_frame
        tf.clear()
        for line in (msg.content or '').split('\n'):
            p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
            p.text = line
            p.font.size = Pt(14)
    out_path = Path(output_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def pptx_to_markdown(input_path: str | Path) -> str:
    """Extract slide text into Markdown headings/blocks."""
    _require_pptx()
    prs = Presentation(str(input_path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = None
        texts: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = (shape.text or "").strip()
            if not text:
                continue
            if shape == slide.shapes.title:
                title = text
            else:
                texts.append(text)
        parts.append(f"## Slide {i}: {title or 'Untitled'}")
        for text in texts:
            parts.append(text)
        parts.append("")
    return "\n\n".join(parts).strip() + "\n"
